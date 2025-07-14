from flask import Flask, request, render_template, send_file, flash, redirect, url_for
import os
import re
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.secret_key = 'your-secret-key-change-this'  # Change this to a random secret key
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Create uploads directory if it doesn't exist
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

class SlideGenerator:
    def __init__(self):
        self.themes = {
            "Professional Blue": {
                'bg': RGBColor(255, 255, 255),  # White background
                'title': RGBColor(31, 78, 121),  # Dark blue
                'text': RGBColor(68, 68, 68),   # Dark gray
                'accent': RGBColor(79, 129, 189) # Light blue
            },
            "Clean Gray": {
                'bg': RGBColor(250, 250, 250),  # Light gray background
                'title': RGBColor(64, 64, 64),   # Dark gray
                'text': RGBColor(96, 96, 96),   # Medium gray
                'accent': RGBColor(128, 128, 128) # Gray accent
            },
            "Medical Green": {
                'bg': RGBColor(255, 255, 255),  # White background
                'title': RGBColor(46, 125, 50),  # Dark green
                'text': RGBColor(68, 68, 68),   # Dark gray
                'accent': RGBColor(102, 187, 106) # Light green
            }
        }
    
    def parse_slide_content(self, text):
        """Parse the formatted text into slide data"""
        slides = []
        
        # Split by slide numbers (0., 1., 2., etc.)
        slide_pattern = r'(\d+)\.\s+(.+?)(?=\n\d+\.\s+|$)'
        slide_matches = re.findall(slide_pattern, text, re.DOTALL)
        
        for slide_num, content in slide_matches:
            slide_data = {
                'number': int(slide_num),
                'title': '',
                'bullets': [],
                'speaker_notes': ''
            }
            
            # Extract title (first line after slide number)
            lines = content.strip().split('\n')
            if lines:
                slide_data['title'] = lines[0].strip()
            
            # Find slide bullets section
            bullets_start = None
            speaker_notes_start = None
            
            for i, line in enumerate(lines):
                if 'Slide Bullets:' in line:
                    bullets_start = i + 1
                elif 'Speaker Notes:' in line:
                    speaker_notes_start = i + 1
                    break
            
            # Extract bullets
            if bullets_start and speaker_notes_start:
                for i in range(bullets_start, speaker_notes_start - 1):
                    if i < len(lines):
                        bullet = lines[i].strip()
                        if bullet and not bullet.startswith('Speaker Notes:'):
                            slide_data['bullets'].append(bullet)
            
            # Extract speaker notes
            if speaker_notes_start:
                notes_lines = lines[speaker_notes_start:]
                slide_data['speaker_notes'] = '\n'.join(notes_lines).strip()
            
            slides.append(slide_data)
        
        return slides
    
    def create_slide(self, prs, slide_data, colors):
        """Create a single slide with the given data"""
        # Use title and content layout
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data['title']
        title.text_frame.paragraphs[0].font.color.rgb = colors['title']
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Add bullets if content placeholder exists
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()
            
            for i, bullet in enumerate(slide_data['bullets']):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet
                p.font.color.rgb = colors['text']
                p.font.size = Pt(18)
                p.level = 0
        
        # Add speaker notes
        if slide_data['speaker_notes']:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = slide_data['speaker_notes']
    
    def generate_presentation(self, text_content, theme_name):
        """Generate PowerPoint presentation from text content"""
        slides = self.parse_slide_content(text_content)
        
        if not slides:
            raise ValueError("No slides found in the input text")
        
        # Create presentation
        prs = Presentation()
        colors = self.themes.get(theme_name, self.themes["Professional Blue"])
        
        # Create slides
        for slide_data in slides:
            self.create_slide(prs, slide_data, colors)
        
        # Save to memory
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output, len(slides)

# Initialize slide generator
slide_gen = SlideGenerator()

@app.route('/')
def index():
    return render_template('index.html', themes=list(slide_gen.themes.keys()))

@app.route('/generate', methods=['POST'])
def generate_slides():
    try:
        # Check if file was uploaded
        if 'file' not in request.files:
            flash('No file selected')
            return redirect(url_for('index'))
        
        file = request.files['file']
        theme = request.form.get('theme', 'Professional Blue')
        
        if file.filename == '':
            flash('No file selected')
            return redirect(url_for('index'))
        
        if file:
            # Read file content
            content = file.read().decode('utf-8')
            
            # Generate presentation
            output, slide_count = slide_gen.generate_presentation(content, theme)
            
            # Generate filename
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f'slides_{timestamp}.pptx'
            
            # Return file for download
            return send_file(
                output,
                as_attachment=True,
                download_name=filename,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
    
    except Exception as e:
        flash(f'Error generating slides: {str(e)}')
        return redirect(url_for('index'))

@app.route('/health')
def health_check():
    return {'status': 'healthy', 'timestamp': datetime.now().isoformat()}

if __name__ == '__main__':
    # Run the app
    # For development: app.run(debug=True, host='0.0.0.0', port=5000)
    # For production: use a proper WSGI server like gunicorn
    app.run(debug=False, host='0.0.0.0', port=5000)
